import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import io

class NinjaWayPresentation:
    def __init__(self):
        self.prs = Presentation()
        # Set slide size to 16:9 widescreen
        self.prs.slide_width = Inches(16)
        self.prs.slide_height = Inches(9)
        
        # Define color scheme
        self.colors = {
            'naruto_orange': RGBColor(255, 140, 0),
            'naruto_blue': RGBColor(0, 100, 200),
            'black': RGBColor(0, 0, 0),
            'white': RGBColor(255, 255, 255),
            'neon_blue': RGBColor(0, 255, 255),
            'neon_purple': RGBColor(138, 43, 226),
            'neon_pink': RGBColor(255, 20, 147)
        }
        
    def create_slide_1(self):
        """Create the title slide: Our Nigga Way"""
        # Add blank slide
        blank_slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_slide_layout)
        
        # Add background image
        bg_path = "Screenshots/slide1background.png"
        if os.path.exists(bg_path):
            # Add background
            left = top = 0
            pic = slide.shapes.add_picture(bg_path, left, top, 
                                           width=self.prs.slide_width,
                                           height=self.prs.slide_height)
            # Send to back
            slide.shapes._spTree.remove(pic._element)
            slide.shapes._spTree.insert(2, pic._element)
        
        # Add orange/blue overlay effect (semi-transparent rectangle)
        overlay = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, 
            self.prs.slide_width, self.prs.slide_height
        )
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = self.colors['naruto_orange']
        overlay.fill.transparency = 0.85  # Make it very transparent
        overlay.line.fill.background()
        
        # Add main title "Our Nigga Way"
        title_box = slide.shapes.add_textbox(
            Inches(1), Inches(2.5), Inches(14), Inches(2.5)
        )
        title_frame = title_box.text_frame
        title_frame.margin_bottom = Inches(0)
        title_frame.margin_top = Inches(0)
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        title = title_frame.add_paragraph()
        title.text = "OUR NIGGA WAY"
        title.font.size = Pt(96)
        title.font.bold = True
        title.font.name = "Impact"
        title.font.color.rgb = self.colors['naruto_orange']
        title.alignment = PP_ALIGN.CENTER
        
        # Add text shadow effect (by creating duplicate text)
        shadow_box = slide.shapes.add_textbox(
            Inches(1.05), Inches(2.55), Inches(14), Inches(2.5)
        )
        shadow_frame = shadow_box.text_frame
        shadow_frame.margin_bottom = Inches(0)
        shadow_frame.margin_top = Inches(0)
        shadow_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        shadow = shadow_frame.add_paragraph()
        shadow.text = "OUR NIGGA WAY"
        shadow.font.size = Pt(96)
        shadow.font.bold = True
        shadow.font.name = "Impact"
        shadow.font.color.rgb = self.colors['black']
        shadow.alignment = PP_ALIGN.CENTER
        
        # Move shadow behind main text
        slide.shapes._spTree.remove(shadow_box._element)
        slide.shapes._spTree.insert(3, shadow_box._element)
        
        # Add subtitle "Light Skins Unite" with neon effect
        subtitle_box = slide.shapes.add_textbox(
            Inches(3), Inches(5.5), Inches(10), Inches(1.5)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        subtitle = subtitle_frame.add_paragraph()
        subtitle.text = "Light Skins Unite"
        subtitle.font.size = Pt(48)
        subtitle.font.bold = True
        subtitle.font.name = "Arial Black"
        subtitle.font.color.rgb = self.colors['neon_blue']
        subtitle.alignment = PP_ALIGN.CENTER
        
        # Add decorative elements (speed lines effect using shapes)
        for i in range(5):
            line = slide.shapes.add_connector(
                1, Inches(8), Inches(4.5),  # Start from center
                Inches(2 + i*3), Inches(1 + i*0.5)  # End points
            )
            line.line.color.rgb = self.colors['naruto_orange']
            line.line.width = Pt(3)
            line.line.transparency = 0.7
            
        # Add corner bursts
        for corner in [(0.5, 0.5), (15, 0.5), (0.5, 8), (15, 8)]:
            burst = slide.shapes.add_shape(
                MSO_SHAPE.STAR_16_POINT, 
                Inches(corner[0]), Inches(corner[1]),
                Inches(1), Inches(1)
            )
            burst.fill.solid()
            burst.fill.fore_color.rgb = self.colors['naruto_blue']
            burst.fill.transparency = 0.6
            burst.line.fill.background()
            
    def save_presentation(self, filename="our_nigga_way.pptx"):
        """Save the presentation"""
        output_path = os.path.join("secure_powerpoints", filename)
        self.prs.save(output_path)
        print(f"Presentation saved to: {output_path}")
        
    def build_full_presentation(self):
        """Build all slides"""
        self.create_slide_1()
        # We'll add other slides as we design them
        self.save_presentation()

if __name__ == "__main__":
    presentation = NinjaWayPresentation()
    presentation.build_full_presentation() 